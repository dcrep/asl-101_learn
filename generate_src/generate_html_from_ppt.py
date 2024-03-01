import string
import re
from pathlib import Path
from pptx import Presentation

# Dependency:
#  https://python-pptx.readthedocs.io/en/latest/
#  pip install python-pptx
#
# Attributions: scanny @ StackOverflow (for gathering links from Powerpoint slides)
# "python - Extract hyperlink from pptx - Stack Overflow" @ https://stackoverflow.com/a/67217499

# TO-DO:
# - Perhaps enable Embedding for videos/animations, have it pop-up/close when clicked?
# - Multiple columns (perhaps with space for embed area)
# - Include more text from slides - context, practice, techniques, mnemonics, etc
# - Hand-edit to show surrounding context (especially important when ONE word of many is hyperlinked)
#   Also good to incorporate practice-sentences using it
#
# - OR! Translate entire PPTX into HTML, and maybe take out some things

# PROBLEM: Presentation() chokes on malformed slides in .PPTX
# (or at least "malformed" according to what Q/A comments suggest)
# Solutions: Run Powerpoint fix (File-Info-Inspect/Check), or edit it, OR shift slides around,
# OR delete and reconstruct slides
# (for one slide, I copied the slide from another file and edited it to match the one in current pptx)

# Function:
# Iterate through slides in powerpoint file, grabe & write out HTML links in list form
# @path: powerpoint file path
# @file: file handle for writing HTML output to
def write_slide(path, file):	
	# Presentation iteration code derived from scanny's post on StackOverflow
	# "python - Extract hyperlink from pptx - Stack Overflow" @ https://stackoverflow.com/a/67217499
	prs = Presentation(path)

	for slide in prs.slides:
		for shape in slide.shapes:
			if not shape.has_text_frame:
				continue
			for paragraph in shape.text_frame.paragraphs:
				for run in paragraph.runs:
					# hyperlink = run.hyperlink
					# print(hyperlink)
					address = run.hyperlink.address
					if address is None:
						continue
					run_text = run.text
					if run_text is None:
						run_text = address
					# Clean up links with trailing #... sequences:
					address = re.sub(r"#[^\"]+", "", address)
					file.write("<li> <a href=\"" + address + "\">" + run_text + "</a> </li>\n")
					print(run_text)
					print(address)

	


# directory name
dirname = '.'

file = open("output.html", "w")
file.write("<!DOCTYPE html>\n<html lang=\"en\">\n<head>\n\t<meta charset=\"utf-8\">\n")
file.write("\t<title>ASL [Week x] Powerpoint Links &amp; Practice (Christine Chynoweth)</title>\n</head>\n")
file.write("<h1>ASL-101 [Week x] Powerpoint Links &amp; Practice (Christine Chynoweth)</h1>\n")
file.write("<a href=\"./index.html\">Back to Index</a>\n")
file.write("<body>\n")
 
# giving directory name to Path() function
paths = Path(dirname).glob('**/*.pptx',)
 
# iterating over all files
for path in paths:
	# pull out Unit-Day# designation
	lesson_name = re.search(r"U\d+D\d+", str(path)).group(0)
	print("[ = = = " + str(path) + " = = = ]")
	file.write("<h2>" + str(path) + "</h2>\n")
	file.write("<p>\n")
	# each list gets id of Unit-Day#; also adjust line-height
	file.write("<ul id=\"" + lesson_name + "\" style=\"line-height:125%\">\n")
	write_slide(path, file)
	file.write("</ul>\n")
	file.write("</p>\n")

# Resources
file.write("<h2>Resources</h2>\n")

# Overall
file.write("<p>\n")
file.write("<ul>\n")
file.write("<li> <a href=\"" + "https://www.lifeprint.com/" + "\">" + "ASLU (Bill Vicars) - LifePrint Resource" + "</a> </li>\n")
file.write("<li> <a href=\"" + "https://www.handspeak.com/" + "\">" + "HandSpeak ASL Resource" + "</a> </li>\n")
file.write("</ul>\n")
file.write("</p>\n")

file.write("<h3>Numbers</h3>\n")

# Numbers
file.write("<p>\n")
file.write("<a href=\"" + "https://www.youtube.com/watch?v=hFCXyB6q2nU" + "\">" + "Numbers 1 to 30 - ASL" + "</a>\n")
file.write("</p>\n")

# Numbers Charts
file.write("<p>\n")
file.write("<a href=\"" + "https://www.lifeprint.com/asl101/pages-signs/n/numbers1-10.htm" + "\">" + "Numbers 1-10" + "</a>\n")
# Also: https://www.youtube.com/watch?v=x0OA2V35g8g "Numbers 1 to 10 | ASL - American Sign Language"
file.write("&nbsp;&nbsp;")
file.write("<a href=\"" + "https://www.lifeprint.com/asl101/pages-signs/n/numbers11-20.htm" + "\">" + "Numbers 11-20" + "</a>\n")
# Also: https://www.youtube.com/watch?v=ItEahhHutzM "Numbers 11-15 in ASL | American Sign Language for Beginners"
file.write("&nbsp;&nbsp;")
file.write("<a href=\"" + "https://www.lifeprint.com/asl101/pages-signs/n/numbers21-30.htm" + "\">" + "Numbers 21-30" + "</a>\n")
# Also: https://www.youtube.com/watch?v=hFCXyB6q2nU "Numbers 1-30 in ASL | American Sign Language for Beginners"
file.write("</p>\n")

# Numbers - Testing
file.write("<p>\n")
file.write("<a href=\"" + "https://asl.bz/" + "\">" + "Number Testing" + "</a>\n")
file.write("</p>\n")

file.write("<h3>Letters/Spelling</h3>\n")

# Letters/Spelling
file.write("<p>\n")
file.write("<a href=\"" + "https://www.youtube.com/watch?v=tkMg8g8vVUo" + "\">" + "Fingerspelling - ASL Alphabet" + "</a>\n")
file.write("&nbsp;&nbsp;&nbsp;&nbsp;")
file.write("<a href=\"" + "https://www.lifeprint.com/asl101/topics/wallpaper1.htm" + "\">" + "Fingerspelling Chart" + "</a>\n")
file.write("&nbsp;&nbsp;&nbsp;&nbsp;")
file.write("<a href=\"" + "https://asl.gs/" + "\">" + "Fingerspelling (on hover)" + "</a>\n")
file.write("</p>\n")

# Letters/Spelling (cont'd)
file.write("<p>\n")
file.write("<h1>\n")
# Letter GIFs from https://asl.gs/slideshow/a.gif -> https://asl.gs/slideshow/z.gif
for letter in string.ascii_lowercase:
	file.write(" <a href=\"" + "https://asl.gs/slideshow/" + letter + ".gif" + "\">" + letter.upper() + "</a>\n")
	
file.write("</h1>\n")
file.write("</p>\n")

# Letters/Words/Spelling - Testing
file.write("<p>\n")
file.write("<a href=\"" + "https://asl.ms/" + "\">" + "Word/Spell Testing" + "</a>\n")
file.write("&nbsp;&nbsp;&nbsp;&nbsp;")
file.write("<a href=\"" + "https://www.handspeak.com/learn/413/" + "\">" + "Word/Spell Testing - Handspeak Receptive" + "</a>\n")
file.write("</p>\n")

file.write("<a href=\"./index.html\">Back to Index</a>\n")
file.write("</body>\n</html>")
file.close()